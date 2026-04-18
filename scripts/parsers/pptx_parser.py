"""
pptx_parser.py — Extract text from PowerPoint presentations (.pptx) into markdown.

This parser:
  1. Reads the .pptx file.
  2. Extracts slide-by-slide text content including speaker notes.
  3. Outputs one .md file per presentation into parsed/docs/.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu  # noqa: F401 — imported for potential future use


def parse_pptx(
    pptx_path: Path,
    output_dir: Path,
) -> Path:
    """Parse a PowerPoint presentation and write a markdown file.

    Args:
        pptx_path: Path to the .pptx file.
        output_dir: Directory to write the parsed markdown (e.g., parsed/docs/).

    Returns:
        Path to the generated markdown file.
    """
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f"{pptx_path.stem}.md"

    prs = Presentation(str(pptx_path))

    lines: list[str] = []
    title = pptx_path.stem.replace("_", " ").replace("-", " ").title()
    lines.append(f"# {title}")
    lines.append("")
    lines.append(f"> Parsed from `{pptx_path.name}`")
    lines.append(f"> Slides: {len(prs.slides)}")
    lines.append("")
    lines.append("---")
    lines.append("")

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_title = None
        body_texts: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text.strip()
            if not text:
                continue

            # The first shape with text that looks like a title
            if shape.shape_id == slide.shapes.title and slide_title is None:
                slide_title = text
            elif slide_title is None and shape == slide.shapes[0]:
                slide_title = text
            else:
                # Process paragraphs to preserve bullet structure
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if not para_text:
                        continue
                    if para.level and para.level > 0:
                        indent = "  " * para.level
                        body_texts.append(f"{indent}- {para_text}")
                    else:
                        body_texts.append(para_text)

        # Use slide title or fallback
        heading = slide_title or f"Slide {slide_num}"
        lines.append(f"## Slide {slide_num}: {heading}")
        lines.append("")

        if body_texts:
            for text in body_texts:
                lines.append(text)
            lines.append("")

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                lines.append("**Speaker Notes:**")
                lines.append("")
                lines.append(f"> {notes_text}")
                lines.append("")

        lines.append("---")
        lines.append("")

    if len(prs.slides) == 0:
        lines.append("*No slides found in this presentation.*")
        lines.append("")

    output_path.write_text("\n".join(lines), encoding="utf-8")
    return output_path
